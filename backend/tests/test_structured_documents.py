from __future__ import annotations

from io import BytesIO

import fitz
from docx import Document
from openpyxl import Workbook
from openpyxl.worksheet.table import Table
from pptx import Presentation

from codeatlas.documents import (
    StructuredBlock,
    chunk_document,
    extract_structured_blocks,
    split_structured_blocks,
)


def test_markdown_uses_heading_hierarchy_before_size_fallback() -> None:
    blocks = [
        StructuredBlock(
            kind="section",
            title="Overview",
            text="Architecture overview.",
            hierarchy=("Guide", "Overview"),
            ordinal=1,
        ),
        StructuredBlock(
            kind="section",
            title="Operations",
            text="First paragraph.\n\nSecond paragraph.",
            hierarchy=("Guide", "Operations"),
            ordinal=2,
        ),
    ]

    chunks = split_structured_blocks("Project guide", blocks, max_chars=80)

    assert [chunk.section for chunk in chunks] == ["Guide > Overview", "Guide > Operations"]
    assert chunks[0].structure_type == "section"
    assert "Architecture overview." in chunks[0].content
    assert chunks[1].metadata["ordinal"] == 2


def test_oversized_structure_splits_on_paragraphs_not_raw_vector_distance() -> None:
    block = StructuredBlock(
        kind="section",
        title="Large section",
        text="Alpha sentence.\n\nBeta sentence.\n\nGamma sentence.",
        hierarchy=("Large section",),
        ordinal=1,
    )

    chunks = split_structured_blocks("Guide", [block], max_chars=38)

    assert len(chunks) >= 2
    assert all(not chunk.content.endswith("Beta sent") for chunk in chunks)
    assert all(chunk.metadata["split_strategy"] == "paragraph" for chunk in chunks)


def test_docx_preserves_heading_hierarchy_and_tables() -> None:
    document = Document()
    document.add_heading("Architecture", level=1)
    document.add_paragraph("Service overview.")
    document.add_heading("API", level=2)
    document.add_paragraph("REST endpoints.")
    table = document.add_table(rows=2, cols=2)
    table.cell(0, 0).text = "Method"
    table.cell(0, 1).text = "Path"
    table.cell(1, 0).text = "GET"
    table.cell(1, 1).text = "/health"
    stream = BytesIO()
    document.save(stream)

    blocks = extract_structured_blocks("architecture.docx", stream.getvalue())

    assert any(block.hierarchy == ("Architecture", "API") for block in blocks)
    assert any(block.kind == "table" and "Method | Path" in block.text for block in blocks)


def test_xlsx_preserves_sheet_table_and_header_context() -> None:
    workbook = Workbook()
    sheet = workbook.active
    sheet.title = "API Owners"
    sheet.append(["Service", "Owner"])
    sheet.append(["Search", "Platform"])
    sheet.append(["Index", "Data"])
    stream = BytesIO()
    workbook.save(stream)

    blocks = extract_structured_blocks("owners.xlsx", stream.getvalue())

    assert blocks[0].kind == "table"
    assert blocks[0].hierarchy == ("API Owners",)
    assert blocks[0].metadata["header"] == ["Service", "Owner"]
    assert "Search | Platform" in blocks[0].text


def test_xlsx_native_table_preserves_real_source_rows_and_table_name() -> None:
    workbook = Workbook()
    sheet = workbook.active
    sheet.title = "SLA矩阵"
    sheet["A16"] = "Service"
    sheet["B16"] = "Target"
    sheet["A17"] = "Order"
    sheet["B17"] = "800ms"
    sheet["A18"] = "Inventory"
    sheet["B18"] = "450ms"
    sheet.add_table(Table(displayName="SlaMatrix", ref="A16:B18"))
    stream = BytesIO()
    workbook.save(stream)

    blocks = extract_structured_blocks("sla.xlsx", stream.getvalue())

    assert blocks[0].hierarchy == ("SLA矩阵", "SlaMatrix")
    assert blocks[0].metadata["table"] == "SlaMatrix"
    assert blocks[0].metadata["row_start"] == 17
    assert blocks[0].metadata["row_end"] == 18
    assert "Order | 800ms" in blocks[0].text


def test_xlsx_preserves_populated_cells_outside_native_tables() -> None:
    workbook = Workbook()
    sheet = workbook.active
    sheet.title = "SLA矩阵"
    sheet["D2"] = "说明"
    sheet["E2"] = "当前数据窗口为8月"
    sheet["A16"] = "Service"
    sheet["B16"] = "Target"
    sheet["A17"] = "Order"
    sheet["B17"] = "800ms"
    sheet.add_table(Table(displayName="SlaMatrix", ref="A16:B17"))
    stream = BytesIO()
    workbook.save(stream)

    blocks = extract_structured_blocks("sla.xlsx", stream.getvalue())

    table = next(block for block in blocks if block.metadata.get("table") == "SlaMatrix")
    outside = next(block for block in blocks if block.kind == "sheet-cells")
    assert table.metadata["row_start"] == 17
    assert outside.metadata["sheet"] == "SLA矩阵"
    assert outside.metadata["row_start"] == 2
    assert outside.metadata["row_end"] == 2
    assert "D2: 说明" in outside.text
    assert "E2: 当前数据窗口为8月" in outside.text


def test_pdf_preserves_page_numbers_and_detects_scans() -> None:
    document = fitz.open()
    first = document.new_page()
    first.insert_text((72, 72), "Architecture\nService overview")
    document.new_page()
    payload = document.tobytes()

    blocks = extract_structured_blocks("architecture.pdf", payload)

    assert any(block.page == 1 and "Architecture" in block.text for block in blocks)
    assert any(block.kind == "ocr-required" and block.page == 2 for block in blocks)


def test_pptx_preserves_slide_title_body_table_and_notes() -> None:
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide.shapes.title.text = "Roadmap"
    slide.placeholders[1].text = "Index documents"
    table = slide.shapes.add_table(2, 2, 0, 0, 2_000_000, 1_000_000).table
    table.cell(0, 0).text = "Phase"
    table.cell(0, 1).text = "Goal"
    table.cell(1, 0).text = "1"
    table.cell(1, 1).text = "Parse"
    slide.notes_slide.notes_text_frame.text = "Presenter context"
    stream = BytesIO()
    presentation.save(stream)

    blocks = extract_structured_blocks("roadmap.pptx", stream.getvalue())

    assert any(
        block.kind == "slide"
        and block.page == 1
        and "Index documents" in block.text
        for block in blocks
    )
    assert any(block.kind == "table" and "Phase | Goal" in block.text for block in blocks)
    assert any(block.kind == "notes" and "Presenter context" in block.text for block in blocks)


def test_document_chunks_persist_structure_metadata() -> None:
    blocks = [
        StructuredBlock(
            kind="table",
            title="Owners",
            text="Service | Owner\nSearch | Platform",
            hierarchy=("Workbook", "Owners"),
            ordinal=3,
            page=2,
            metadata={"sheet": "Owners", "row_start": 2, "row_end": 3},
        )
    ]

    chunks = chunk_document("Owner matrix", "doc-1", "collection-1", blocks=blocks)

    assert chunks[0].structure_type == "table"
    assert chunks[0].page == 2
    metadata = __import__("json").loads(chunks[0].metadata_json)
    assert metadata["sheet"] == "Owners"
    assert metadata["row_start"] == 2
