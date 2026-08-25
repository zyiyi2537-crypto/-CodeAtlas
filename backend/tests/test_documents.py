from __future__ import annotations

from io import BytesIO
from pathlib import Path

import pymupdf
from fastapi.testclient import TestClient
from pptx import Presentation
from sqlmodel import Session

from codeatlas.models import Document, DocumentChunkRecord, DocumentCollection
from tests.conftest import login_admin


def test_admin_can_upload_markdown_document(
    client: TestClient, admin, tmp_path: Path
) -> None:
    csrf = login_admin(client)
    headers = {"X-CSRF-Token": csrf}
    collection = client.post(
        "/api/v1/document-collections",
        headers=headers,
        json={"name": "订单文档", "description": "订单项目资料"},
    )
    assert collection.status_code == 201
    collection_id = collection.json()["id"]

    uploaded = client.post(
        f"/api/v1/document-collections/{collection_id}/documents",
        headers=headers,
        files={
            "file": (
                "refund-standard.md",
                b"# Refund\n\n## Idempotency\nRefunds must be idempotent.",
                "text/markdown",
            )
        },
        data={"title": "退款规范"},
    )
    assert uploaded.status_code == 201
    payload = uploaded.json()
    assert payload["title"] == "退款规范"
    assert payload["status"] == "indexed"
    assert payload["version"] == 1

    search = client.post(
        "/api/v1/documents/search",
        headers=headers,
        json={"query": "退款 幂等", "collection_ids": [collection_id]},
    )
    assert search.status_code == 200
    assert search.json()[0]["source_type"] == "document"
    assert "idempotent" in search.json()[0]["content"]


def test_document_upload_rejects_unsupported_type(client: TestClient, admin) -> None:
    csrf = login_admin(client)
    headers = {"X-CSRF-Token": csrf}
    collection = client.post(
        "/api/v1/document-collections",
        headers=headers,
        json={"name": "资料"},
    )
    response = client.post(
        f"/api/v1/document-collections/{collection.json()['id']}/documents",
        headers=headers,
        files={"file": ("bad.exe", b"MZ", "application/octet-stream")},
    )
    assert response.status_code == 415


def test_index_failed_document_is_not_searchable(application, admin, monkeypatch) -> None:
    with Session(application.state.engine) as session:
        collection = DocumentCollection(
            name="Failed documents", description="", created_by=admin.id
        )
        session.add(collection)
        session.flush()
        collection_id = collection.id
        document = Document(
            collection_id=collection.id,
            title="Failed deployment",
            original_filename="failed.md",
            mime_type="text/markdown",
            status="index_failed",
            source_path="",
            sha256="0" * 64,
            created_by=admin.id,
        )
        session.add(document)
        session.flush()
        document_id = document.id
        session.add(
            DocumentChunkRecord(
                id="failed-document-chunk",
                document_id=document.id,
                collection_id=collection.id,
                title=document.title,
                section="Secrets",
                content="This failed document must not appear in retrieval.",
            )
        )
        session.commit()

    monkeypatch.setattr(
        application.state.knowledge_search.vector_store,
        "search_knowledge",
        lambda *_args, **_kwargs: [
            {
                "id": "failed-document-chunk",
                "document": "This failed document must not appear in retrieval.",
                "metadata": {
                    "source_type": "document",
                    "source_id": document_id,
                    "collection_id": collection_id,
                    "title": "Failed deployment",
                    "section": "Secrets",
                    "page": 0,
                },
                "vector_score": 0.99,
            }
        ],
    )

    results = application.state.knowledge_search.search_documents("failed document")

    assert results == []


def test_admin_can_upload_text_pdf_and_search_by_page(client: TestClient, admin) -> None:
    csrf = login_admin(client)
    headers = {"X-CSRF-Token": csrf}
    collection = client.post(
        "/api/v1/document-collections",
        headers=headers,
        json={"name": "PDF 文档"},
    ).json()
    pdf = pymupdf.open()
    page = pdf.new_page()
    page.insert_text((72, 72), "Deployment\nConfigure the reverse proxy")

    uploaded = client.post(
        f"/api/v1/document-collections/{collection['id']}/documents",
        headers=headers,
        files={"file": ("deployment.pdf", pdf.tobytes(), "application/pdf")},
    )

    assert uploaded.status_code == 201
    search = client.post(
        "/api/v1/documents/search",
        headers=headers,
        json={"query": "reverse proxy", "collection_ids": [collection["id"]]},
    )
    assert search.status_code == 200
    assert search.json()[0]["page"] == 1
    assert search.json()[0]["retrieval"] in {"hybrid", "vector"}


def test_scanned_pdf_is_marked_for_ocr_without_indexing_placeholder(
    client: TestClient, admin
) -> None:
    csrf = login_admin(client)
    headers = {"X-CSRF-Token": csrf}
    collection = client.post(
        "/api/v1/document-collections",
        headers=headers,
        json={"name": "扫描件"},
    ).json()
    pdf = pymupdf.open()
    pdf.new_page()

    uploaded = client.post(
        f"/api/v1/document-collections/{collection['id']}/documents",
        headers=headers,
        files={"file": ("scan.pdf", pdf.tobytes(), "application/pdf")},
    )

    assert uploaded.status_code == 201
    assert uploaded.json()["status"] == "ocr_required"
    assert uploaded.json()["chunk_count"] == 0


def test_admin_can_upload_pptx_and_retrieve_slide(client: TestClient, admin) -> None:
    csrf = login_admin(client)
    headers = {"X-CSRF-Token": csrf}
    collection = client.post(
        "/api/v1/document-collections",
        headers=headers,
        json={"name": "PPT 文档"},
    ).json()
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide.shapes.title.text = "Roadmap"
    slide.placeholders[1].text = "Build unified semantic retrieval"
    stream = BytesIO()
    presentation.save(stream)

    uploaded = client.post(
        f"/api/v1/document-collections/{collection['id']}/documents",
        headers=headers,
        files={
            "file": (
                "roadmap.pptx",
                stream.getvalue(),
                "application/vnd.openxmlformats-officedocument.presentationml.presentation",
            )
        },
    )

    assert uploaded.status_code == 201
    search = client.post(
        "/api/v1/documents/search",
        headers=headers,
        json={"query": "semantic retrieval", "collection_ids": [collection["id"]]},
    )
    assert search.status_code == 200
    assert search.json()[0]["page"] == 1


def test_unified_knowledge_search_returns_document_results(client: TestClient, admin) -> None:
    csrf = login_admin(client)
    headers = {"X-CSRF-Token": csrf}
    collection = client.post(
        "/api/v1/document-collections",
        headers=headers,
        json={"name": "统一检索"},
    ).json()
    uploaded = client.post(
        f"/api/v1/document-collections/{collection['id']}/documents",
        headers=headers,
        files={
            "file": (
                "guide.md",
                b"# Deployment\n\nConfigure Nginx reverse proxy.",
                "text/markdown",
            )
        },
    )
    assert uploaded.status_code == 201

    response = client.post(
        "/api/v1/knowledge/search",
        headers=headers,
        json={
            "query": "Nginx reverse proxy",
            "source_types": ["code", "document", "wiki"],
            "collection_ids": [collection["id"]],
        },
    )

    assert response.status_code == 200
    assert response.json()[0]["source_type"] == "document"
