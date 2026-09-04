from __future__ import annotations

import hashlib
import json
import re
from dataclasses import dataclass, field
from io import BytesIO
from pathlib import Path
from typing import Any

import pymupdf
from docx import Document as DocxDocument
from docx.table import Table as DocxTable
from docx.text.paragraph import Paragraph as DocxParagraph
from openpyxl import load_workbook
from openpyxl.utils.cell import range_boundaries
from pptx import Presentation

from .models import DEFAULT_SPACE_ID, DocumentChunkRecord
from .security import redact_secrets

_ALLOWED = {".md", ".markdown", ".txt", ".csv", ".docx", ".xlsx", ".pdf", ".pptx"}


@dataclass(frozen=True)
class StructuredBlock:
    kind: str
    title: str
    text: str
    hierarchy: tuple[str, ...] = ()
    ordinal: int = 0
    page: int | None = None
    metadata: dict[str, Any] = field(default_factory=dict)


@dataclass(frozen=True)
class StructuredChunk:
    section: str
    structure_type: str
    content: str
    page: int | None
    metadata: dict[str, Any]


def _paragraph_parts(text: str, max_chars: int) -> list[str]:
    paragraphs = [item.strip() for item in re.split(r"\n\s*\n", text) if item.strip()]
    parts: list[str] = []
    current = ""
    for paragraph in paragraphs:
        candidate = f"{current}\n\n{paragraph}".strip() if current else paragraph
        if current and len(candidate) > max_chars:
            parts.append(current)
            current = paragraph
        else:
            current = candidate
        while len(current) > max_chars:
            boundaries = [
                current.rfind("。", 0, max_chars),
                current.rfind(". ", 0, max_chars),
                current.rfind("；", 0, max_chars),
                current.rfind("; ", 0, max_chars),
            ]
            boundary = max(boundaries)
            boundary = boundary + 1 if boundary >= max_chars // 2 else max_chars
            parts.append(current[:boundary].strip())
            current = current[boundary:].strip()
    if current:
        parts.append(current)
    return parts


def split_structured_blocks(
    document_title: str, blocks: list[StructuredBlock], max_chars: int = 3500
) -> list[StructuredChunk]:
    chunks: list[StructuredChunk] = []
    for block in blocks:
        if block.kind == "ocr-required":
            continue
        hierarchy = block.hierarchy or ((block.title,) if block.title else ())
        section = " > ".join(item for item in hierarchy if item)
        for index, part in enumerate(_paragraph_parts(block.text, max_chars), start=1):
            content = (
                f"Document: {document_title}\nSection: {section}\n"
                f"Structure: {block.kind}\n\n{part}"
            )
            chunks.append(
                StructuredChunk(
                    section=section,
                    structure_type=block.kind,
                    content=content,
                    page=block.page,
                    metadata={
                        **block.metadata,
                        "ordinal": block.ordinal,
                        "part": index,
                        "split_strategy": "paragraph",
                    },
                )
            )
    return chunks


def _markdown_blocks(text: str) -> list[StructuredBlock]:
    blocks: list[StructuredBlock] = []
    hierarchy: list[str] = []
    body: list[str] = []
    ordinal = 0

    def flush() -> None:
        nonlocal ordinal, body
        value = "\n".join(body).strip()
        if value:
            ordinal += 1
            blocks.append(
                StructuredBlock(
                    kind="section",
                    title=hierarchy[-1] if hierarchy else "",
                    text=value,
                    hierarchy=tuple(hierarchy),
                    ordinal=ordinal,
                )
            )
        body = []

    for line in text.splitlines():
        match = re.match(r"^(#{1,6})\s+(.+)$", line.strip())
        if match:
            flush()
            level = len(match.group(1))
            hierarchy[:] = hierarchy[: level - 1]
            hierarchy.append(match.group(2).strip())
        else:
            body.append(line)
    flush()
    return blocks


def _docx_blocks(content: bytes) -> list[StructuredBlock]:
    document = DocxDocument(BytesIO(content))
    blocks: list[StructuredBlock] = []
    hierarchy: list[str] = []
    ordinal = 0
    body: list[str] = []

    def flush_body() -> None:
        nonlocal ordinal, body
        text = "\n\n".join(item for item in body if item).strip()
        if text:
            ordinal += 1
            blocks.append(
                StructuredBlock(
                    kind="section",
                    title=hierarchy[-1] if hierarchy else "",
                    text=text,
                    hierarchy=tuple(hierarchy),
                    ordinal=ordinal,
                )
            )
        body = []

    for item in document.iter_inner_content():
        if isinstance(item, DocxParagraph):
            text = item.text.strip()
            if not text:
                continue
            style = item.style.name if item.style else ""
            match = re.match(r"Heading\s+([1-6])$", style, re.IGNORECASE)
            if match:
                flush_body()
                level = int(match.group(1))
                hierarchy[:] = hierarchy[: level - 1]
                hierarchy.append(text)
            else:
                body.append(text)
        elif isinstance(item, DocxTable):
            flush_body()
            rows = [
                [cell.text.strip().replace("\n", " ") for cell in row.cells]
                for row in item.rows
            ]
            if rows:
                ordinal += 1
                blocks.append(
                    StructuredBlock(
                        kind="table",
                        title=hierarchy[-1] if hierarchy else "Table",
                        text="\n".join(" | ".join(row) for row in rows),
                        hierarchy=tuple(hierarchy),
                        ordinal=ordinal,
                        metadata={"header": rows[0], "row_count": len(rows) - 1},
                    )
                )
    flush_body()
    return blocks


def _xlsx_table_blocks(
    sheet, table_name: str, reference: str, ordinal: int, rows_per_block: int
) -> tuple[list[StructuredBlock], int]:
    min_column, min_row, max_column, max_row = range_boundaries(reference)
    rows = [
        (
            row_number,
            [
                "" if sheet.cell(row_number, column).value is None
                else str(sheet.cell(row_number, column).value)
                for column in range(min_column, max_column + 1)
            ],
        )
        for row_number in range(min_row, max_row + 1)
    ]
    if not rows:
        return [], ordinal
    header = rows[0][1]
    data = [(row_number, row) for row_number, row in rows[1:] if any(row)]
    groups = [
        data[offset : offset + rows_per_block]
        for offset in range(0, len(data), rows_per_block)
    ] or [[]]
    blocks = []
    for group in groups:
        ordinal += 1
        row_start = group[0][0] if group else min_row + 1
        row_end = group[-1][0] if group else min_row
        blocks.append(
            StructuredBlock(
                kind="table",
                title=table_name,
                text="\n".join(
                    [" | ".join(header), *[" | ".join(row) for _, row in group]]
                ),
                hierarchy=(sheet.title, table_name),
                ordinal=ordinal,
                metadata={
                    "sheet": sheet.title,
                    "table": table_name,
                    "header": header,
                    "row_start": row_start,
                    "row_end": row_end,
                },
            )
        )
    return blocks, ordinal


def _xlsx_outside_table_blocks(
    sheet,
    covered_ranges: list[tuple[int, int, int, int]],
    ordinal: int,
    rows_per_block: int,
) -> tuple[list[StructuredBlock], int]:
    def covered(row: int, column: int) -> bool:
        return any(
            min_row <= row <= max_row and min_column <= column <= max_column
            for min_column, min_row, max_column, max_row in covered_ranges
        )

    populated_rows: list[tuple[int, list[str]]] = []
    for row_number in range(1, sheet.max_row + 1):
        cells = []
        for column in range(1, sheet.max_column + 1):
            if covered(row_number, column):
                continue
            value = sheet.cell(row_number, column).value
            if value is None or str(value).strip() == "":
                continue
            coordinate = sheet.cell(row_number, column).coordinate
            cells.append(f"{coordinate}: {value}")
        if cells:
            populated_rows.append((row_number, cells))

    groups: list[list[tuple[int, list[str]]]] = []
    current: list[tuple[int, list[str]]] = []
    for row in populated_rows:
        if current and (row[0] != current[-1][0] + 1 or len(current) >= rows_per_block):
            groups.append(current)
            current = []
        current.append(row)
    if current:
        groups.append(current)

    blocks = []
    for group in groups:
        ordinal += 1
        blocks.append(
            StructuredBlock(
                kind="sheet-cells",
                title=sheet.title,
                text="\n".join(" | ".join(cells) for _, cells in group),
                hierarchy=(sheet.title, "Cells outside tables"),
                ordinal=ordinal,
                metadata={
                    "sheet": sheet.title,
                    "row_start": group[0][0],
                    "row_end": group[-1][0],
                },
            )
        )
    return blocks, ordinal


def _xlsx_blocks(content: bytes, rows_per_block: int = 50) -> list[StructuredBlock]:
    workbook = load_workbook(BytesIO(content), data_only=False, read_only=False)
    blocks: list[StructuredBlock] = []
    ordinal = 0
    for sheet in workbook.worksheets:
        tables = sorted(sheet.tables.values(), key=lambda table: table.ref)
        if tables:
            covered_ranges = []
            for table in tables:
                covered_ranges.append(range_boundaries(table.ref))
                table_blocks, ordinal = _xlsx_table_blocks(
                    sheet,
                    table.displayName,
                    table.ref,
                    ordinal,
                    rows_per_block,
                )
                blocks.extend(table_blocks)
            outside_blocks, ordinal = _xlsx_outside_table_blocks(
                sheet,
                covered_ranges,
                ordinal,
                rows_per_block,
            )
            blocks.extend(outside_blocks)
            continue
        rows = [
            (
                row_number,
                ["" if value is None else str(value) for value in row],
            )
            for row_number, row in enumerate(sheet.iter_rows(values_only=True), start=1)
        ]
        rows = [(row_number, row) for row_number, row in rows if any(cell.strip() for cell in row)]
        if not rows:
            continue
        header, data = rows[0][1], rows[1:]
        groups = [
            data[offset : offset + rows_per_block]
            for offset in range(0, len(data), rows_per_block)
        ] or [[]]
        for group in groups:
            ordinal += 1
            blocks.append(
                StructuredBlock(
                    kind="table",
                    title=sheet.title,
                    text="\n".join(
                        [" | ".join(header), *[" | ".join(row) for _, row in group]]
                    ),
                    hierarchy=(sheet.title,),
                    ordinal=ordinal,
                    metadata={
                        "sheet": sheet.title,
                        "header": header,
                        "row_start": group[0][0] if group else rows[0][0] + 1,
                        "row_end": group[-1][0] if group else rows[0][0],
                    },
                )
            )
    return blocks


def _pdf_blocks(content: bytes) -> list[StructuredBlock]:
    document = pymupdf.open(stream=content, filetype="pdf")
    blocks: list[StructuredBlock] = []
    for page_number in range(1, document.page_count + 1):
        page = document.load_page(page_number - 1)
        raw_blocks = sorted(page.get_text("blocks"), key=lambda item: (round(item[1]), item[0]))
        texts = [str(item[4]).strip() for item in raw_blocks if str(item[4]).strip()]
        if not texts:
            blocks.append(
                StructuredBlock(
                    kind="ocr-required",
                    title=f"Page {page_number}",
                    text=f"Page {page_number} contains no extractable text and requires OCR.",
                    hierarchy=(f"Page {page_number}",),
                    ordinal=page_number,
                    page=page_number,
                    metadata={"requires_ocr": True},
                )
            )
            continue
        title = next(
            (line for line in texts[0].splitlines() if line.strip()),
            f"Page {page_number}",
        )
        blocks.append(
            StructuredBlock(
                kind="page",
                title=title,
                text="\n\n".join(texts),
                hierarchy=(title,),
                ordinal=page_number,
                page=page_number,
                metadata={"requires_ocr": False},
            )
        )
    return blocks


def _pptx_blocks(content: bytes) -> list[StructuredBlock]:
    presentation = Presentation(BytesIO(content))
    blocks: list[StructuredBlock] = []
    ordinal = 0
    for slide_number, slide in enumerate(presentation.slides, start=1):
        title_shape = slide.shapes.title
        title = title_shape.text.strip() if title_shape else f"Slide {slide_number}"
        body = []
        for shape in slide.shapes:
            if shape is title_shape or not getattr(shape, "has_text_frame", False):
                continue
            text = shape.text.strip()
            if text:
                body.append(text)
        if body:
            ordinal += 1
            blocks.append(
                StructuredBlock(
                    kind="slide",
                    title=title,
                    text="\n\n".join(body),
                    hierarchy=(title,),
                    ordinal=ordinal,
                    page=slide_number,
                    metadata={"slide": slide_number},
                )
            )
        for shape in slide.shapes:
            if not getattr(shape, "has_table", False):
                continue
            rows = [
                [cell.text.strip().replace("\n", " ") for cell in row.cells]
                for row in shape.table.rows
            ]
            ordinal += 1
            blocks.append(
                StructuredBlock(
                    kind="table",
                    title=title,
                    text="\n".join(" | ".join(row) for row in rows),
                    hierarchy=(title,),
                    ordinal=ordinal,
                    page=slide_number,
                    metadata={"slide": slide_number, "header": rows[0] if rows else []},
                )
            )
        notes = slide.notes_slide.notes_text_frame.text.strip()
        notes = "\n".join(
            line for line in notes.splitlines() if line.strip() and line.strip() != title
        ).strip()
        if notes:
            ordinal += 1
            blocks.append(
                StructuredBlock(
                    kind="notes",
                    title=title,
                    text=notes,
                    hierarchy=(title, "Notes"),
                    ordinal=ordinal,
                    page=slide_number,
                    metadata={"slide": slide_number},
                )
            )
    return blocks


def extract_structured_blocks(filename: str, content: bytes) -> list[StructuredBlock]:
    suffix = Path(filename).suffix.lower()
    if suffix not in _ALLOWED:
        raise ValueError(
            "supported document types are Markdown, TXT, CSV, DOCX, XLSX, PDF and PPTX"
        )
    try:
        if suffix == ".docx":
            blocks = _docx_blocks(content)
        elif suffix == ".xlsx":
            blocks = _xlsx_blocks(content)
        elif suffix == ".pdf":
            blocks = _pdf_blocks(content)
        elif suffix == ".pptx":
            blocks = _pptx_blocks(content)
        else:
            blocks = _markdown_blocks(content.decode("utf-8-sig"))
    except UnicodeDecodeError as exc:
        raise ValueError("document must be UTF-8 text") from exc
    except (KeyError, ValueError, OSError) as exc:
        raise ValueError(f"invalid {suffix.removeprefix('.').upper()} document") from exc
    return [
        StructuredBlock(
            kind=block.kind,
            title=block.title,
            text=redact_secrets(block.text),
            hierarchy=block.hierarchy,
            ordinal=block.ordinal,
            page=block.page,
            metadata=block.metadata,
        )
        for block in blocks
    ]


def extract_text(filename: str, content: bytes) -> str:
    return "\n\n".join(block.text for block in extract_structured_blocks(filename, content))


def chunk_document(
    title: str,
    document_id: str,
    collection_id: str,
    text: str | None = None,
    max_chars: int = 3500,
    *,
    blocks: list[StructuredBlock] | None = None,
    space_id: str = DEFAULT_SPACE_ID,
) -> list[DocumentChunkRecord]:
    structured = blocks if blocks is not None else _markdown_blocks(text or "")
    chunks: list[DocumentChunkRecord] = []
    for chunk in split_structured_blocks(title, structured, max_chars):
        metadata_json = json.dumps(
            {
                "structure_type": chunk.structure_type,
                **chunk.metadata,
            },
            ensure_ascii=False,
        )
        chunk_id = hashlib.sha256(
            f"{document_id}|{chunk.section}|{chunk.page}|{metadata_json}|{chunk.content}".encode()
        ).hexdigest()
        chunks.append(
            DocumentChunkRecord(
                id=chunk_id,
                document_id=document_id,
                collection_id=collection_id,
                space_id=space_id,
                title=title,
                section=chunk.section,
                page=chunk.page,
                structure_type=chunk.structure_type,
                metadata_json=metadata_json,
                content=chunk.content,
            )
        )
    return chunks
